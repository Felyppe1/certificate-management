from threading import Event, Lock, current_thread
import functions_framework
from dotenv import load_dotenv
import os
from google.oauth2.credentials import Credentials
import requests
from io import BytesIO, StringIO
from docx import Document
from docx.text.paragraph import Paragraph
from pptx import Presentation
from google.cloud import storage
from google.cloud.storage.blob import Blob
import re
import google.auth.transport.requests
import google.oauth2.id_token
from liquid import RenderContext
from liquid_types import LiquidDate, LiquidFloat, liquid_environment
from pydantic import BaseModel, ValidationError
from typing import List, Optional, Dict, Any
from enum import Enum
from datetime import datetime
import re
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload

load_dotenv()

ENV = os.getenv('ENV', 'production')
APP_BASE_URL = os.getenv('APP_BASE_URL')
AUDIENCE = os.getenv("TOKEN_AUDIENCE", APP_BASE_URL) # For prod, it will use the APP_BASE_URL as AUDIENCE
CERTIFICATES_BUCKET = os.getenv('CERTIFICATES_BUCKET')
GOOGLE_DRIVE_FOLDER_ID = os.getenv('GOOGLE_DRIVE_FOLDER_ID')
GOOGLE_REFRESH_TOKEN = os.getenv('GOOGLE_REFRESH_TOKEN')
GOOGLE_CLIENT_ID = os.getenv('GOOGLE_CLIENT_ID')
GOOGLE_CLIENT_SECRET = os.getenv('GOOGLE_CLIENT_SECRET')

for var_name, var_value in {
    "APP_BASE_URL": APP_BASE_URL,
    "CERTIFICATES_BUCKET": CERTIFICATES_BUCKET,
    "GOOGLE_DRIVE_FOLDER_ID": GOOGLE_DRIVE_FOLDER_ID,
    "GOOGLE_REFRESH_TOKEN": GOOGLE_REFRESH_TOKEN,
    "GOOGLE_CLIENT_ID": GOOGLE_CLIENT_ID,
    "GOOGLE_CLIENT_SECRET": GOOGLE_CLIENT_SECRET
}.items():
    if not var_value:
        raise ValueError(f"Environment variable '{var_name}' is not set.")

DOCX_MIME_TYPE = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
GOOGLE_DOCS_MIME_TYPE = 'application/vnd.google-apps.document'
PPTX_MIME_TYPE = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
GOOGLE_SLIDES_MIME_TYPE = 'application/vnd.google-apps.presentation'

storage_client = storage.Client()

BLOCK_START_TAGS = {
    'if', 'unless', 'case', 'for', 'tablerow', 'capture', 'form', 'paginate', 'raw', 'comment', 'liquid'
}
BLOCK_END_TAGS = {
    'endif', 'endunless', 'endcase', 'endfor', 'endtablerow', 'endcapture', 'endform', 'endpaginate', 'endraw', 'endcomment', 'endliquid'
}

# Regex to find opening tags: {% if, {%- if, {% endfor, etc.
TAG_REGEX = re.compile(r'{%-?\s*(\w+)')

#################################### Functions to process DOCX with Liquid ####################################
def paragraph_universal_iterator(doc):
    """Returns all body, table, and textbox paragraphs."""
    for p in doc.paragraphs:
        yield p
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    yield p
    
    ns_map = doc.element.nsmap
    txbx_contents = doc.element.body.findall('.//w:txbxContent', ns_map)
    for txbx in txbx_contents:
        xml_paragraphs = txbx.findall('.//w:p', ns_map)
        for xml_p in xml_paragraphs:
            yield Paragraph(xml_p, doc)

def consolidate_broken_runs(paragraph):
    """Merges runs within the SAME paragraph to fix Word breaks."""
    text = paragraph.text
    if not ('{{' in text or '{%' in text):
        return

    runs = paragraph.runs
    if not runs: return

    i = 0
    while i < len(runs):
        text_run = runs[i].text
        has_open = '{%' in text_run or '{{' in text_run
        
        if has_open:
            closure = '%}' if '{%' in text_run else '}}'
            j = i + 1
            while closure not in runs[i].text and j < len(runs):
                next_run = runs[j]
                runs[i].text += next_run.text
                next_run.text = ""
                j += 1
        i += 1

def calculate_delta_blocks(text):
    """
    Returns the variation of opened/closed blocks in this text.
    +1 for each 'if', 'for'...
    -1 for each 'endif', 'endfor'...
    """
    delta = 0
    # Finds all tags {% tag ... %}
    matches = TAG_REGEX.findall(text)
    
    for tag_name in matches:
        if tag_name in BLOCK_START_TAGS:
            delta += 1
        elif tag_name in BLOCK_END_TAGS:
            delta -= 1
    return delta

def process_buffer(paragraph_list, text_list, liquid_ctx):
    '''
    Takes N paragraphs, joins the text, runs Liquid, and returns the result
    in the FIRST paragraph, deleting the others.
    '''
    # Joins with line breaks to simulate the original document
    complete_text = "\n".join(text_list)

    # If there is no liquid syntax, ignore
    if not ("{{" in complete_text or "{%" in complete_text):
        return

    sanitized_text = (complete_text
                      .replace('“', '"').replace('”', '"')
                      .replace('‘', "'").replace('’', "'"))
    if "append" in sanitized_text:
        print(f"DEBUG LIQUID: {sanitized_text}")

    try:
        template = liquid_environment.from_string(sanitized_text)
        buf = StringIO()
        template.render_with_context(liquid_ctx, buf)
        new_complete_text = buf.getvalue()
        # If the text did not change, do nothing (preserves original formatting)
        if new_complete_text == complete_text:
            return

        # MULTI-LINE REPLACEMENT STRATEGY:
        # 1. Put EVERYTHING in the first paragraph of the buffer
        main_paragraph = paragraph_list[0]
        
        # Clear runs of the main paragraph
        for run in main_paragraph.runs:
            run.text = ""
        # Add the new text (may contain line breaks resulting from liquid)
        # Note: python-docx handles '\n' inside a run by creating soft breaks, 
        # but visually it works for the purpose.
        if main_paragraph.runs:
            main_paragraph.runs[0].text = new_complete_text
        else:
            main_paragraph.add_run(new_complete_text)

        # 2. Clear the subsequent paragraphs that were part of the block
        # Ex: The paragraph that had the "{% else %}" and the "{% endif %}"
        for p_lixo in paragraph_list[1:]:
            for run in p_lixo.runs:
                run.text = ""
            # Optional: If you want to remove the paragraph entirely (may break layout if it's a table)
            # p_lixo._element.getparent().remove(p_lixo._element) 
            # Recommended to just clear the text to keep the table/doc structure safe.

    except Exception as e:
        print(f"Error Liquid block started in '{text_list[0][:20]}...': {e}")

def replace_variables_in_docx_liquid(in_buffer, context_data):
    in_buffer.seek(0)
    doc = Document(in_buffer)

    # 1. Pre-processing: Consolidate runs (intra-line)
    # This ensures that '{% if' is not split across different runs in the same line
    all_paragraphs = list(paragraph_universal_iterator(doc))
    for p in all_paragraphs:
        consolidate_broken_runs(p)

    # Shared context so {% assign %} variables persist across paragraph renders
    liquid_ctx = RenderContext(liquid_environment.from_string(""), globals=context_data)

    # 2. Processing with Buffer (Inter-line)
    buffer_paragraphs = []
    accumulated_text = []
    block_level = 0

    idx = 0
    while idx < len(all_paragraphs):
        p = all_paragraphs[idx]
        paragraph_text = p.text

        delta = calculate_delta_blocks(paragraph_text)

        if block_level > 0 or (delta > 0):
            buffer_paragraphs.append(p)
            accumulated_text.append(paragraph_text)
            block_level += delta

            if block_level == 0:
                process_buffer(buffer_paragraphs, accumulated_text, liquid_ctx)
                buffer_paragraphs = []
                accumulated_text = []

        else:
            process_buffer([p], [paragraph_text], liquid_ctx)

        idx += 1

    if buffer_paragraphs:
        print("WARNING: Liquid block not closed at the end of the document.")
        process_buffer(buffer_paragraphs, accumulated_text, liquid_ctx)

    # 3. Save
    out_buffer = BytesIO()
    doc.save(out_buffer)
    out_buffer.seek(0)
    return out_buffer


#################################### Functions to process PPTX with Liquid ####################################
def pptx_paragraph_iterator(prs):
    """
    Yields tuples of (paragraph, shape) for all paragraphs in the presentation.
    Includes text frames, tables, and grouped shapes.
    """
    def process_shape(shape):
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                yield (paragraph, shape)

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        yield (paragraph, shape)

        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for subshape in shape.shapes:
                yield from process_shape(subshape)

    for slide in prs.slides:
        for shape in slide.shapes:
            yield from process_shape(shape)


def consolidate_broken_runs_pptx(paragraph):
    """Merges runs within the SAME paragraph to fix PowerPoint breaks."""
    text = "".join(run.text for run in paragraph.runs)
    if not ('{{' in text or '{%' in text):
        return

    runs = paragraph.runs
    if not runs:
        return

    i = 0
    while i < len(runs):
        text_run = runs[i].text
        has_open = '{%' in text_run or '{{' in text_run
        
        if has_open:
            closure = '%}' if '{%' in text_run else '}}'
            j = i + 1
            while closure not in runs[i].text and j < len(runs):
                next_run = runs[j]
                runs[i].text += next_run.text
                next_run.text = ""
                j += 1
        i += 1


def process_buffer_pptx(paragraph_list, text_list, liquid_ctx):
    """
    Takes N paragraphs, joins the text, runs Liquid, and returns the result
    in the FIRST paragraph, clearing the others.
    """
    complete_text = "\n".join(text_list)

    if not ("{{" in complete_text or "{%" in complete_text):
        return

    sanitized_text = (complete_text
                      .replace('“', '"').replace('”', '"')
                      .replace('‘', "'").replace('’', "'"))

    if "append" in sanitized_text:
        print(f"DEBUG LIQUID PPTX: {sanitized_text}")

    try:
        template = liquid_environment.from_string(sanitized_text)
        buf = StringIO()
        template.render_with_context(liquid_ctx, buf)
        new_complete_text = buf.getvalue()
        
        if new_complete_text == complete_text:
            return

        main_paragraph = paragraph_list[0]
        
        for run in main_paragraph.runs:
            run.text = ""
        
        if main_paragraph.runs:
            main_paragraph.runs[0].text = new_complete_text
        else:
            main_paragraph.add_run().text = new_complete_text

        for p in paragraph_list[1:]:
            for run in p.runs:
                run.text = ""

    except Exception as e:
        print(f"Error Liquid block started in '{text_list[0][:20]}...': {e}")


def replace_variables_in_pptx_liquid(in_buffer, context_data):
    in_buffer.seek(0)
    prs = Presentation(in_buffer)

    # 1. Pre-processing: Consolidate runs (intra-line)
    all_paragraphs = list(pptx_paragraph_iterator(prs))
    for p, _ in all_paragraphs:
        consolidate_broken_runs_pptx(p)

    # Shared context so {% assign %} variables persist across paragraph renders
    liquid_ctx = RenderContext(liquid_environment.from_string(""), globals=context_data)

    # 2. Processing with Buffer (Inter-line)
    buffer_paragraphs = []
    accumulated_text = []
    block_level = 0

    idx = 0
    while idx < len(all_paragraphs):
        p, _ = all_paragraphs[idx]
        paragraph_text = "".join(run.text for run in p.runs)

        delta = calculate_delta_blocks(paragraph_text)

        if block_level > 0 or (delta > 0):
            buffer_paragraphs.append(p)
            accumulated_text.append(paragraph_text)
            block_level += delta

            if block_level == 0:
                process_buffer_pptx(buffer_paragraphs, accumulated_text, liquid_ctx)
                buffer_paragraphs = []
                accumulated_text = []

        else:
            process_buffer_pptx([p], [paragraph_text], liquid_ctx)

        idx += 1

    if buffer_paragraphs:
        print("WARNING: Liquid block not closed at the end of the presentation.")
        process_buffer_pptx(buffer_paragraphs, accumulated_text, liquid_ctx)

    # 3. Save
    out_buffer = BytesIO()
    prs.save(out_buffer)
    out_buffer.seek(0)
    return out_buffer


#################################### Function to convert certificate to pdf ####################################
def convert_to_pdf_with_google_drive(input_bytes: BytesIO, input_ext: str) -> BytesIO:
    """
    Convert DOCX or PPTX to PDF using Google Drive API:
    1. Upload the file to the Google Drive
    2. Export it as PDF
    3. Delete the original file from Drive
    """
    mime_types = {
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    }
    
    export_mime_types = {
        'docx': 'application/vnd.google-apps.document',
        'pptx': 'application/vnd.google-apps.presentation'
    }
    
    mime_type = mime_types.get(input_ext)
    if not mime_type:
        raise ValueError(f"Unsupported file extension: {input_ext}")
    
    credentials = Credentials(
        None,
        refresh_token=GOOGLE_REFRESH_TOKEN,
        token_uri="https://oauth2.googleapis.com/token",
        client_id=GOOGLE_CLIENT_ID,
        client_secret=GOOGLE_CLIENT_SECRET,
        scopes=["https://www.googleapis.com/auth/drive.file", 'https://www.googleapis.com/auth/drive.readonly',],
    )
    drive_service = build('drive', 'v3', credentials=credentials)
    
    try:
        file_metadata = {
            'name': f'temp_convert.{input_ext}',
            'parents': [GOOGLE_DRIVE_FOLDER_ID],
            'mimeType': export_mime_types[input_ext]
        }
        
        input_bytes.seek(0)
        media = MediaIoBaseUpload(input_bytes, mimetype=mime_type, resumable=False)
        
        uploaded_file = drive_service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id'
        ).execute()
        
        file_id = uploaded_file.get('id')
        print(f"File sent to the Drive with ID: {file_id}")
        
        pdf_content = drive_service.files().export(
            fileId=file_id,
            mimeType='application/pdf'
        ).execute()
        
        pdf_buffer = BytesIO(pdf_content)
        pdf_buffer.seek(0)
        
        drive_service.files().delete(fileId=file_id).execute()
        print(f"File {file_id} deleted from Drive")
        
        return pdf_buffer
        
    except Exception as e:
        print(f"Error converting with Google Drive: {e}")
        raise


#################################### Functions to use bucket ####################################
def upload_to_bucket(file_buffer, file_path, content_type="application/pdf"):
    bucket = storage_client.bucket(CERTIFICATES_BUCKET)
    blob = bucket.blob(file_path)
    blob.upload_from_file(file_buffer, rewind=True, content_type=content_type)
    return blob

def get_from_bucket(file_path) -> Blob:
    bucket = storage_client.bucket(CERTIFICATES_BUCKET)
    blob = bucket.blob(file_path)
    return blob

CACHE = {}
IN_FLIGHT = {}
LOCK = Lock()

def get_template_cached(storage_file_url: str) -> bytes:
    thread_name = current_thread().name
    print(f"[{thread_name}] Requesting template: {storage_file_url}")

    template_blob = get_from_bucket(storage_file_url)
    template_blob.reload()
    generation = str(template_blob.generation)

    with LOCK:
        cached = CACHE.get(storage_file_url)
        if cached and cached["generation"] == generation:
            print(f"[{thread_name}] CACHE HIT (generation={generation})")
            return cached["bytes"]

        print(f"[{thread_name}] CACHE MISS")

        if storage_file_url in IN_FLIGHT:
            print(f"[{thread_name}] Another thread is downloading → waiting")
            event = IN_FLIGHT[storage_file_url]
            is_downloader = False
        else:
            print(f"[{thread_name}] Elected as downloader")
            event = Event()
            IN_FLIGHT[storage_file_url] = event
            is_downloader = True

    # 🔹 Threads that DO NOT download wait here
    if not is_downloader:
        event.wait()
        print(f"[{thread_name}] Download finished → reading from cache")
        return CACHE[storage_file_url]["bytes"]

    # 🔹 Only ONE thread arrives here
    try:
        print(f"[{thread_name}] Downloading template (generation={generation})")
        template_bytes = template_blob.download_as_bytes()

        with LOCK:
            CACHE[storage_file_url] = {
                "generation": generation,
                "bytes": template_bytes
            }
            print(f"[{thread_name}] Cache updated")
    finally:
        with LOCK:
            event.set()
            IN_FLIGHT.pop(storage_file_url, None)
            print(f"[{thread_name}] Released waiting threads")

    return template_bytes


#################################### Functions to call backend endpoints ####################################
def finish_certificates_generation(data_source_row_id, success, total_bytes=None, user_id=None):
    print('Inside update')
    url = f"{APP_BASE_URL}/api/internal/data-source-rows/{data_source_row_id}/generations"

    if ENV != 'local':
        auth_req = google.auth.transport.requests.Request()
        # Just works for service accounts or in the cloud. Locally, I need to login impersonating a service account.
        id_token = google.oauth2.id_token.fetch_id_token(auth_req, AUDIENCE)

        headers = {
            "Authorization": f"Bearer {id_token}",
            "Content-Type": "application/json",
        }
    else:
        headers = {
            "Content-Type": "application/json",
        }
    
    body = {k: v for k, v in {
        "success": success,
        "totalBytes": total_bytes,
        "userId": user_id,
    }.items() if v is not None}

    print('before sending patch')
    response = requests.patch(url, json=body, headers=headers)
    response.raise_for_status()



#################################### Input schemas ####################################
class InputMethod(str, Enum):
    UPLOAD = "UPLOAD"
    URL = "URL"
    GOOGLE_DRIVE = "GOOGLE_DRIVE"

class TemplateFileMimeType(str, Enum):
    PPTX = 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    GOOGLE_SLIDES = 'application/vnd.google-apps.presentation',
    GOOGLE_DOCS = 'application/vnd.google-apps.document',
    DOCX = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',

class DataSourceFileExtension(str, Enum):
    CSV = 'text/csv',
    XLSX = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    ODS = 'application/vnd.oasis.opendocument.spreadsheet',
    GOOGLE_SHEETS = 'application/vnd.google-apps.spreadsheet'

class DataSourceRowModel(BaseModel):
    id: str
    data: Dict[str, Any]

class TemplateModel(BaseModel):
    storageFileUrl: Optional[str] = None
    fileMimeType: TemplateFileMimeType
    variables: List[str]

class ColumnType(str, Enum):
    STRING = "string"
    NUMBER = "number"
    BOOLEAN = "boolean"
    DATE = "date"
    ARRAY = "array"

class ArrayItemType(str, Enum):
    STRING = "string"
    NUMBER = "number"
    BOOLEAN = "boolean"
    DATE = "date"

class ArrayMetadata(BaseModel):
    separator: str
    itemType: ArrayItemType

class Column(BaseModel):
    name: str
    type: ColumnType
    arrayMetadata: Optional[ArrayMetadata] = None

class DataSourceModel(BaseModel):
    columns: List[Column]

class CertificateEmissionModel(BaseModel):
    id: str
    userId: str
    variableColumnMapping: Optional[Dict[str, Optional[str]]] = None # Record<string, string | null> | null
    template: TemplateModel
    dataSource: DataSourceModel

class TriggerGenerateCertificatePDFsInput(BaseModel):
    certificateEmission: CertificateEmissionModel
    row: DataSourceRowModel

@functions_framework.http
def main(request):
    print('Generate PDFs function invoked via Pub/Sub Push')
    
    raw_data = request.get_json(silent=True)
    if raw_data is None:
        return {"error": "JSON body is required"}, 400

    data_source_row_id = raw_data.get('row', {}).get('id')
    user_id = None
    
    try:
        input_data = TriggerGenerateCertificatePDFsInput(**raw_data)
        
        certificate_emission = input_data.certificateEmission
        certificate_emission_id = certificate_emission.id
        template = certificate_emission.template
        data_source = certificate_emission.dataSource
        variable_mapping = certificate_emission.variableColumnMapping or {}
        row = input_data.row
        data_source_row_id = row.id
        user_id = certificate_emission.userId
        # input_method = template.inputMethod.value
        file_mime_type = template.fileMimeType.value

        template_buffer = None
        is_docx = None

        if file_mime_type in [DOCX_MIME_TYPE, GOOGLE_DOCS_MIME_TYPE]:
            is_docx = True
            file_extension_str = 'docx'
        elif file_mime_type in [PPTX_MIME_TYPE, GOOGLE_SLIDES_MIME_TYPE]:
            is_docx = False
            file_extension_str = 'pptx'
        else:
            raise Exception(f'Unsupported template file extension: {file_mime_type}')

        print('Loading template from bucket: ', template.storageFileUrl)
        template_bytes = get_template_cached(template.storageFileUrl)
        template_buffer = BytesIO(template_bytes)

        print(f'Generating certificate for row {data_source_row_id}: ', row)
        certificate_buffer = BytesIO(template_buffer.getvalue())

        print('variable_mapping: ', variable_mapping)
        if variable_mapping:
            row_variable_mapping = {}
            for template_var, column_name in variable_mapping.items():
                if column_name and column_name in row.data:
                    for column in data_source.columns:
                        if column.name == column_name:
                            def convert_item(item: str, item_type: str):
                                match item_type:
                                    case 'boolean':
                                        return item.lower() in ("true", "verdadeiro", "1")
                                    case 'number':
                                        return LiquidFloat(item)
                                    case 'date':
                                        BR_REGEX = re.compile(
                                            r'^(\d{1,2})/(\d{1,2})/(\d{4})(?: (\d{2}):(\d{2})(?::(\d{2}))?)?$'
                                        )

                                        regex_match = BR_REGEX.match(item)

                                        if regex_match:
                                            day = int(regex_match.group(1))
                                            month = int(regex_match.group(2))
                                            year = int(regex_match.group(3))
                                            hour = int(regex_match.group(4)) if regex_match.group(4) else None
                                            minute = int(regex_match.group(5)) if regex_match.group(5) else None
                                            second = int(regex_match.group(6)) if regex_match.group(6) else None

                                            if month > 12:
                                                aux = day
                                                day = month
                                                month = aux

                                            if second != None:
                                                dt = datetime(year, month, day, hour, minute, second)
                                                return LiquidDate(dt.strftime("%d/%m/%Y %H:%M:%S"), dt)
                                            elif hour != None or minute != None:
                                                dt = datetime(year, month, day, hour, minute)
                                                return LiquidDate(dt.strftime("%d/%m/%Y %H:%M"), dt)
                                            else:
                                                dt = datetime(year, month, day)
                                                return LiquidDate(dt.strftime("%d/%m/%Y"), dt)

                                        return item
                                    case _:
                                        return item

                            row_value = row.data[column_name].strip()
                            
                            if column.type == ColumnType.ARRAY:
                                raw_items = [
                                    item.strip()
                                    for item in row_value.split(column.arrayMetadata.separator)
                                    if item.strip()
                                ]
                                item_type = column.arrayMetadata.itemType

                                row_variable_mapping[template_var] = [convert_item(i, item_type) for i in raw_items]
                            else:
                                row_variable_mapping[template_var] = convert_item(row_value, column.type)

                            break

            if is_docx:
                certificate_buffer = replace_variables_in_docx_liquid(certificate_buffer, row_variable_mapping)
            else:
                certificate_buffer = replace_variables_in_pptx_liquid(certificate_buffer, row_variable_mapping)
        
        source_mime = (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            if is_docx
            else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        source_path = f"users/{user_id}/certificates/{certificate_emission_id}/certificate-{data_source_row_id}.{file_extension_str}"
        upload_to_bucket(certificate_buffer, source_path, content_type=source_mime)

        pdf_buffer = convert_to_pdf_with_google_drive(certificate_buffer, file_extension_str)

        pdf_path = f"users/{user_id}/certificates/{certificate_emission_id}/certificate-{data_source_row_id}.pdf"
        blob = upload_to_bucket(pdf_buffer, pdf_path)

        finish_certificates_generation(data_source_row_id, True, blob.size, user_id)
        
        return "", 204
        
    except ValidationError as e:
        def format_pydantic_errors(errors):
            formatted_errors = []

            for error in errors:
                field = ".".join([str(x) for x in error['loc']])
                message = error['msg']
                formatted_errors.append({
                    "field": field,
                    "message": message
                })

            return formatted_errors

        friendly_errors = format_pydantic_errors(e.errors())

        print("Validation errors:", friendly_errors)
        
        if data_source_row_id and user_id:
            finish_certificates_generation(data_source_row_id, False)
            
        return {"error": friendly_errors}, 200
    
    except Exception as e:
        original_error = str(e)
        update_error = None

        try:
            print('Sending error status update...')
            if data_source_row_id and user_id:
                finish_certificates_generation(data_source_row_id, False)
        except Exception as inner_e:
            update_error = str(inner_e)

        details = original_error if not update_error else f'Original error: {original_error}; Update error: {update_error}'

        print('Error details:', details)
        return {
            'title': 'Failed to generate certificates',
            'details': details
        }, 200